# Please write a Python script that performs the following tasks:

# - Accesses the "input" folder and reads all Excel files within it.

# - use column "Planned Hours" column (B) and "Actual Hours column" (C) .
# - use column "Total Test Cases" column (B) and "Total Test Cases Passed" column (C) . implement that above code

# - The script should create a bar chart, line chart , pie chart,  for each file, and save it in the "charts" folder. If the folder does not exist, the script should create it. . The chart should include a title, appropriate axis labels, and a legend. implement that above code

# - The script should then create a new PowerPoint presentation, insert a slide for each chart, and above the chart, include a title. The title should be the respective excel file name without file extension. And make sure that the chart and title do not overlap
# - The PowerPoint presentation should be saved in the same directory as the input files and named "WSR_MultipleWorksheet.pptx".
# - The script should be robust and handle any potential errors gracefully, providing appropriate error messages and notifications, and also by including proper error handling mechanisms

# WSR_MultipleWorksheet.pptx should be outside of input folder. Implement this

# Charts should be editable in presentation. Implement this
# User able to see only one bar chart.

# Error adding charts to presentation for India_WSR - Project1: 'XySeriesData' object has no attribute 'categories'

# Manage chart area. Currently charts size is not looks good

# Error adding charts to presentation for India_WSR - Project2: 'Planned Hours'
'''import os
import pandas as pd
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_LEGEND_POSITION

# Define the input and output directories
input_folder = 'input'
output_pptx = 'WSR_MultipleWorksheet.pptx'


def add_charts_to_presentation(prs, df, filename):
    try:

        required_columns = ['Planned Hours', 'Actual Hours']

        if not all(col in df.columns for col in required_columns):
            print(f"Error: One or more required columns are missing in {filename}. Skipping this file.")
            return

        # Title layout
        slide_layout = prs.slide_layouts[5]  # Title and Content layout

        # Chart size and position
        x, y, cx, cy = Cm(2.54), Cm(4), Cm(20), Cm(12)

        # Bar Chart
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = f'{filename} - Bar Chart'

        bar_chart_data = CategoryChartData()
        bar_chart_data.categories = df[df.columns[0]]
        bar_chart_data.add_series('Planned Hours', df['Planned Hours'])
        bar_chart_data.add_series('Actual Hours', df['Actual Hours'])

        chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, bar_chart_data).chart

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.has_title = True
        chart.chart_title.text_frame.text = f'{filename} - Planned vs Actual Hours'

        # Line Chart
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = f'{filename} - Line Chart'

        line_chart_data = CategoryChartData()
        line_chart_data.categories = df[df.columns[0]]
        line_chart_data.add_series('Planned Hours', df['Planned Hours'])
        line_chart_data.add_series('Actual Hours', df['Actual Hours'])

        chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, x, y, cx, cy, line_chart_data).chart

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.has_title = True
        chart.chart_title.text_frame.text = f'{filename} - Planned vs Actual Hours'

        # Pie Chart
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = f'{filename} - Pie Chart'

        pie_chart_data = CategoryChartData()
        pie_chart_data.categories = df[df.columns[0]]
        pie_chart_data.add_series('Planned Hours', df['Planned Hours'])

        chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, cx, cy, pie_chart_data).chart

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False
        chart.has_title = True
        chart.chart_title.text_frame.text = f'{filename} - Planned Hours Distribution'

    except Exception as e:
        print(f"Error adding charts to presentation for {filename}: {e}")


def read_excel_files(input_folder):
    excel_files = [f for f in os.listdir(input_folder) if f.endswith('.xlsx') or f.endswith('.xls')]

    prs = Presentation()

    for filename in excel_files:
        file_path = os.path.join(input_folder, filename)
        try:
            # Read the Excel file into a DataFrame
            df = pd.read_excel(file_path, sheet_name=0, usecols='A:C')

            # Drop rows with any missing values
            df.dropna(inplace=True)

            # Add charts to presentation
            base_filename = os.path.splitext(filename)[0]
            add_charts_to_presentation(prs, df, base_filename)

            print(f"Successfully processed {filename}")

        except Exception as e:
            print(f"Error reading {filename}: {e}")

    # Save the PowerPoint presentation outside the input folder
    prs.save(output_pptx)
    print(f"PowerPoint presentation saved to {output_pptx}")


if __name__ == "__main__":
    read_excel_files(input_folder)'''
# merge all chart in slide
'''import os
import pandas as pd
from pptx import Presentation
from pptx.util import Cm
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

# Define the input and output directories
input_folder = 'input'
output_pptx = 'WSR_MultipleWorksheet.pptx'


def add_charts_to_presentation(prs, df, filename):
    try:
        # Title layout
        slide_layout = prs.slide_layouts[5]  # Title and Content layout

        # Chart size and position
        x, y, cx, cy = Cm(2.54), Cm(2), Cm(18), Cm(12)

        # Check if columns exist
        if 'Planned Hours' in df.columns and 'Actual Hours' in df.columns:
            # Create slide for charts
            slide = prs.slides.add_slide(slide_layout)

            # Bar Chart
            bar_chart_data = CategoryChartData()
            bar_chart_data.categories = df[df.columns[0]]
            bar_chart_data.add_series('Planned Hours', df['Planned Hours'])
            bar_chart_data.add_series('Actual Hours', df['Actual Hours'])

            chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, bar_chart_data).chart
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
            chart.has_title = True
            chart.chart_title.text_frame.text = f'{filename} - Planned vs Actual Hours'

            # Line Chart
            x += cx + Cm(1)  # Move right of previous chart
            line_chart_data = CategoryChartData()
            line_chart_data.categories = df[df.columns[0]]
            line_chart_data.add_series('Planned Hours', df['Planned Hours'])
            line_chart_data.add_series('Actual Hours', df['Actual Hours'])

            chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, x, y, cx, cy, line_chart_data).chart
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
            chart.has_title = True
            chart.chart_title.text_frame.text = f'{filename} - Planned vs Actual Hours'

            # Pie Chart
            x += cx + Cm(1)  # Move right of previous chart
            pie_chart_data = CategoryChartData()
            pie_chart_data.categories = df[df.columns[0]]
            pie_chart_data.add_series('Planned Hours', df['Planned Hours'])

            chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, cx, cy, pie_chart_data).chart
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.RIGHT
            chart.legend.include_in_layout = False
            chart.has_title = True
            chart.chart_title.text_frame.text = f'{filename} - Planned Hours Distribution'

        else:
            print(
                f"Error adding charts to presentation for {filename}: 'Planned Hours' or 'Actual Hours' columns not found.")

    except Exception as e:
        print(f"Error adding charts to presentation for {filename}: {e}")


def read_excel_files(input_folder):
    excel_files = [f for f in os.listdir(input_folder) if f.endswith('.xlsx') or f.endswith('.xls')]

    prs = Presentation()

    for filename in excel_files:
        file_path = os.path.join(input_folder, filename)
        try:
            # Read the Excel file into a DataFrame
            df = pd.read_excel(file_path, sheet_name=0, usecols='A:C')

            # Drop rows with any missing values
            df.dropna(inplace=True)

            # Add charts to presentation
            base_filename = os.path.splitext(filename)[0]
            add_charts_to_presentation(prs, df, base_filename)

            print(f"Successfully processed {filename}")

        except Exception as e:
            print(f"Error reading {filename}: {e}")

    # Save the PowerPoint presentation outside the input folder
    prs.save(output_pptx)
    print(f"PowerPoint presentation saved to {output_pptx}")


if __name__ == "__main__":
    read_excel_files(input_folder)'''

'''import os
import pandas as pd
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_LEGEND_POSITION

# Define the input and output directories
input_folder = 'input'
output_pptx = 'WSR_MultipleWorksheet.pptx'


def add_charts_to_presentation(prs, df, filename):
    try:
        required_columns = ['Planned Hours', 'Actual Hours']

        if not all(col in df.columns for col in required_columns):
            print(f"Error: One or more required columns are missing in {filename}. Skipping this file.")
            return

        # Title layout
        slide_layout = prs.slide_layouts[5]  # Title and Content layout

        # Chart size and position
        x, y, cx, cy = Cm(2.54), Cm(4), Cm(20), Cm(12)

        # Bar Chart
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = f'{filename} - Bar Chart'
        title.text_frame.paragraphs[0].font.size = Pt(18)  # Adjust title font size here

        bar_chart_data = CategoryChartData()
        bar_chart_data.categories = df[df.columns[0]]
        bar_chart_data.add_series('Planned Hours', df['Planned Hours'])
        bar_chart_data.add_series('Actual Hours', df['Actual Hours'])

        chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, bar_chart_data).chart

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.has_title = True
        chart.chart_title.text_frame.text = f'{filename} - Planned vs Actual Hours'

        # Line Chart
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = f'{filename} - Line Chart'
        title.text_frame.paragraphs[0].font.size = Pt(18)  # Adjust title font size here

        line_chart_data = CategoryChartData()
        line_chart_data.categories = df[df.columns[0]]
        line_chart_data.add_series('Planned Hours', df['Planned Hours'])
        line_chart_data.add_series('Actual Hours', df['Actual Hours'])

        chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, x, y, cx, cy, line_chart_data).chart

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.has_title = True
        chart.chart_title.text_frame.text = f'{filename} - Planned vs Actual Hours'

        # Pie Chart
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = f'{filename} - Pie Chart'
        title.text_frame.paragraphs[0].font.size = Pt(18)  # Adjust title font size here

        pie_chart_data = CategoryChartData()
        pie_chart_data.categories = df[df.columns[0]]
        pie_chart_data.add_series('Planned Hours', df['Planned Hours'])

        chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, cx, cy, pie_chart_data).chart

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False
        chart.has_title = True
        chart.chart_title.text_frame.text = f'{filename} - Planned Hours Distribution'

    except Exception as e:
        print(f"Error adding charts to presentation for {filename}: {e}")


def read_excel_files(input_folder):
    excel_files = [f for f in os.listdir(input_folder) if f.endswith('.xlsx') or f.endswith('.xls')]

    prs = Presentation()

    for filename in excel_files:
        file_path = os.path.join(input_folder, filename)
        try:
            # Read the Excel file into a DataFrame
            df = pd.read_excel(file_path, sheet_name=0, usecols='A:C')

            # Drop rows with any missing values
            df.dropna(inplace=True)

            # Add charts to presentation
            base_filename = os.path.splitext(filename)[0]
            add_charts_to_presentation(prs, df, base_filename)

            print(f"Successfully processed {filename}")

        except Exception as e:
            print(f"Error reading {filename}: {e}")

    # Save the PowerPoint presentation outside the input folder
    prs.save(output_pptx)
    print(f"PowerPoint presentation saved to {output_pptx}")


if __name__ == "__main__":
    read_excel_files(input_folder)'''

'''import os
import pandas as pd
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_LEGEND_POSITION

# Define the input and output directories
input_folder = 'input'
output_pptx = 'WSR_MultipleWorksheet.pptx'


def add_charts_to_presentation(prs, df, filename):
    try:
        required_columns = ['Planned Hours', 'Actual Hours']

        if not all(col in df.columns for col in required_columns):
            print(f"Error: One or more required columns are missing in {filename}. Skipping this file.")
            return

        # Title layout
        slide_layout = prs.slide_layouts[5]  # Title and Content layout

        # Create a new slide for the charts
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = f'{filename} - Charts'
        title.text_frame.paragraphs[0].font.size = Pt(18)  # Adjust title font size here

        # Chart size and position
        x1, y1, cx1, cy1 = Cm(1), Cm(4), Cm(12), Cm(6)
        x2, y2, cx2, cy2 = Cm(12), Cm(4), Cm(12), Cm(6)
        x3, y3, cx3, cy3 = Cm(1), Cm(11), Cm(12), Cm(6)

        # Bar Chart
        bar_chart_data = CategoryChartData()
        bar_chart_data.categories = df[df.columns[0]]
        bar_chart_data.add_series('Planned Hours', df['Planned Hours'])
        bar_chart_data.add_series('Actual Hours', df['Actual Hours'])

        chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x1, y1, cx1, cy1, bar_chart_data).chart

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.has_title = True
        chart.chart_title.text_frame.text = f'{filename} - Planned vs Actual Hours (Bar Chart)'

        # Line Chart
        line_chart_data = CategoryChartData()
        line_chart_data.categories = df[df.columns[0]]
        line_chart_data.add_series('Planned Hours', df['Planned Hours'])
        line_chart_data.add_series('Actual Hours', df['Actual Hours'])

        chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, x2, y2, cx2, cy2, line_chart_data).chart

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.has_title = True
        chart.chart_title.text_frame.text = f'{filename} - Planned vs Actual Hours (Line Chart)'

        # Pie Chart
        pie_chart_data = CategoryChartData()
        pie_chart_data.categories = df[df.columns[0]]
        pie_chart_data.add_series('Planned Hours', df['Planned Hours'])

        chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, x3, y3, cx3, cy3, pie_chart_data).chart

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False
        chart.has_title = True
        chart.chart_title.text_frame.text = f'{filename} - Planned Hours Distribution (Pie Chart)'

    except Exception as e:
        print(f"Error adding charts to presentation for {filename}: {e}")


def read_excel_files(input_folder):
    excel_files = [f for f in os.listdir(input_folder) if f.endswith('.xlsx') or f.endswith('.xls')]

    prs = Presentation()

    for filename in excel_files:
        file_path = os.path.join(input_folder, filename)
        try:
            # Read the Excel file into a DataFrame
            df = pd.read_excel(file_path, sheet_name=0, usecols='A:C')

            # Drop rows with any missing values
            df.dropna(inplace=True)

            # Add charts to presentation
            base_filename = os.path.splitext(filename)[0]
            add_charts_to_presentation(prs, df, base_filename)

            print(f"Successfully processed {filename}")

        except Exception as e:
            print(f"Error reading {filename}: {e}")

    # Save the PowerPoint presentation outside the input folder
    prs.save(output_pptx)
    print(f"PowerPoint presentation saved to {output_pptx}")


if __name__ == "__main__":
    read_excel_files(input_folder)'''
# provide option to adjust Font size for chart titles ,labels and category axis. Implement in above code
# To provide options to adjust the vertical value axis and legend values of charts, you can add parameters to the add_charts_to_presentation function. Here's how you can modify the function to include these options:
'''import os
import pandas as pd
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_TICK_MARK
from pptx.dml.color import RGBColor

# Define the input and output directories
input_folder = 'input'
output_pptx = 'WSR_MultipleWorksheet.pptx'


def add_charts_to_presentation(prs, df, filename, title_font_size=Pt(12), label_font_size=Pt(8),
                               category_axis_font_size=Pt(8), value_axis_font_size=Pt(8), legend_font_size=Pt(8)):
    try:
        required_columns = ['Planned Hours', 'Actual Hours']

        if not all(col in df.columns for col in required_columns):
            print(f"Error: One or more required columns are missing in {filename}. Skipping this file.")
            return

        # Title layout
        slide_layout = prs.slide_layouts[5]  # Title and Content layout

        # Create a new slide for the charts
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = f'{filename} - Charts'
        title.text_frame.paragraphs[0].font.size = title_font_size  # Adjust title font size here

        # Chart size and position
        gap = Cm(0.5)  # Horizontal gap between charts

        # Chart size and position
        x1, y1, cx1, cy1 = Cm(1), Cm(4), Cm(12), Cm(6)
        x2, y2, cx2, cy2 = Cm(13), Cm(4), Cm(12), Cm(6)
        x3, y3, cx3, cy3 = Cm(1), Cm(11), Cm(12), Cm(6)

        # Bar Chart
        bar_chart_data = CategoryChartData()
        bar_chart_data.categories = df[df.columns[0]]
        bar_chart_data.add_series('Planned Hours', df['Planned Hours'])
        bar_chart_data.add_series('Actual Hours', df['Actual Hours'])

        chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x1, y1, cx1, cy1, bar_chart_data).chart

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = legend_font_size  # Adjust legend font size
        chart.has_title = True
        chart.chart_title.text_frame.text = f'- Planned vs Actual Hours (Bar Chart)'
        chart.chart_title.text_frame.paragraphs[0].font.size = title_font_size  # Adjust chart title font size

        # Adjust font size for labels and category axis
        for series in chart.series:
            for point in series.points:
                point.data_label.font.size = label_font_size

        category_axis = chart.category_axis
        category_axis.tick_labels.font.size = category_axis_font_size

        value_axis = chart.value_axis
        value_axis.tick_labels.font.size = value_axis_font_size

        # Line Chart
        line_chart_data = CategoryChartData()
        line_chart_data.categories = df[df.columns[0]]
        line_chart_data.add_series('Planned Hours', df['Planned Hours'])
        line_chart_data.add_series('Actual Hours', df['Actual Hours'])

        chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, x2, y2, cx2, cy2, line_chart_data).chart

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = legend_font_size  # Adjust legend font size
        chart.has_title = True
        chart.chart_title.text_frame.text = f'- Planned vs Actual Hours (Line Chart)'
        chart.chart_title.text_frame.paragraphs[0].font.size = title_font_size  # Adjust chart title font size

        # Adjust font size for labels and category axis
        for series in chart.series:
            for point in series.points:
                point.data_label.font.size = label_font_size

        category_axis = chart.category_axis
        category_axis.tick_labels.font.size = category_axis_font_size

        value_axis = chart.value_axis
        value_axis.tick_labels.font.size = value_axis_font_size

        # Pie Chart
        pie_chart_data = CategoryChartData()
        pie_chart_data.categories = df[df.columns[0]]
        pie_chart_data.add_series('Planned Hours', df['Planned Hours'])

        chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, x3, y3, cx3, cy3, pie_chart_data).chart

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False
        chart.legend.font.size = legend_font_size  # Adjust legend font size
        chart.has_title = True
        chart.chart_title.text_frame.text = f' - Planned Hours Distribution (Pie Chart)'
        chart.chart_title.text_frame.paragraphs[0].font.size = title_font_size  # Adjust chart title font size

        # Adjust font size for labels
        for series in chart.series:
            for point in series.points:
                point.data_label.font.size = label_font_size

    except Exception as e:
        print(f"Error adding charts to presentation for {filename}: {e}")


def read_excel_files(input_folder):
    excel_files = [f for f in os.listdir(input_folder) if f.endswith('.xlsx') or f.endswith('.xls')]

    prs = Presentation()

    for filename in excel_files:
        file_path = os.path.join(input_folder, filename)
        try:
            # Read the Excel file into a DataFrame
            df = pd.read_excel(file_path, sheet_name=0, usecols='A:C')

            # Drop rows with any missing values
            df.dropna(inplace=True)

            # Add charts to presentation with custom font sizes
            base_filename = os.path.splitext(filename)[0]
            add_charts_to_presentation(prs, df, base_filename, title_font_size=Pt(15), label_font_size=Pt(10),
                                       category_axis_font_size=Pt(10), value_axis_font_size=Pt(10), legend_font_size=Pt(10))

            print(f"Successfully processed {filename}")

        except Exception as e:
            print(f"Error reading {filename}: {e}")

    # Save the PowerPoint presentation outside the input folder
    prs.save(output_pptx)
    print(f"PowerPoint presentation saved to {output_pptx}")


if __name__ == "__main__":
    read_excel_files(input_folder)'''
# To add one more horizontal chart in the above code and adjust the positions accordingly, follow these steps. I will add a fourth chart and ensure all four charts fit well on the slide with proper spacing.
# Explanation:
# Chart Positions and Sizes: Adjusted x and y coordinates, along with widths (cx) and heights (cy) for each chart to ensure they fit on the slide.
# Fourth Chart (Area Chart): Added a fourth chart (Area Chart) with appropriate data series and positioning.
# Chart Titles and Fonts: Ensured all chart titles and fonts are set according to the given parameters.
'''import os
import pandas as pd
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_TICK_MARK
from pptx.dml.color import RGBColor

# Define the input and output directories
input_folder = 'input'
output_pptx = 'WSR_MultipleWorksheet.pptx'


def add_charts_to_presentation(prs, df, filename, title_font_size=Pt(12), label_font_size=Pt(8),
                               category_axis_font_size=Pt(8), value_axis_font_size=Pt(8), legend_font_size=Pt(8)):
    try:
        required_columns = ['Planned Hours', 'Actual Hours']

        if not all(col in df.columns for col in required_columns):
            print(f"Error: One or more required columns are missing in {filename}. Skipping this file.")
            return

        # Title layout
        slide_layout = prs.slide_layouts[5]  # Title and Content layout

        # Create a new slide for the charts
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = f'{filename} - Charts'
        title.text_frame.paragraphs[0].font.size = title_font_size  # Adjust title font size here

        # Chart size and position
        gap = Cm(0.5)  # Horizontal gap between charts

        # Chart size and position # (HP,VP,  W,H)
        x1, y1, cx1, cy1 = Cm(1.5), Cm(4), Cm(9.9), Cm(6)
        x2, y2, cx2, cy2 = x1 + cx1 + gap, Cm(4), Cm(9.9), Cm(6)
        x3, y3, cx3, cy3 = Cm(1.5), y1 + cy1 + gap, Cm(9.9), Cm(6)
        x4, y4, cx4, cy4 = x3 + cx3 + gap, y3, Cm(9.9), Cm(6)

        # Bar Chart
        bar_chart_data = CategoryChartData()
        bar_chart_data.categories = df[df.columns[0]]
        bar_chart_data.add_series('Planned Hours', df['Planned Hours'])
        bar_chart_data.add_series('Actual Hours', df['Actual Hours'])

        chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x1, y1, cx1, cy1, bar_chart_data).chart

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = legend_font_size  # Adjust legend font size
        chart.has_title = True
        chart.chart_title.text_frame.text = f'- Planned vs Actual Hours (Bar Chart)'
        chart.chart_title.text_frame.paragraphs[0].font.size = title_font_size  # Adjust chart title font size

        # Adjust font size for labels and category axis
        for series in chart.series:
            for point in series.points:
                point.data_label.font.size = label_font_size

        category_axis = chart.category_axis
        category_axis.tick_labels.font.size = category_axis_font_size

        value_axis = chart.value_axis
        value_axis.tick_labels.font.size = value_axis_font_size

        # Line Chart
        line_chart_data = CategoryChartData()
        line_chart_data.categories = df[df.columns[0]]
        line_chart_data.add_series('Planned Hours', df['Planned Hours'])
        line_chart_data.add_series('Actual Hours', df['Actual Hours'])

        chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, x2, y2, cx2, cy2, line_chart_data).chart

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = legend_font_size  # Adjust legend font size
        chart.has_title = True
        chart.chart_title.text_frame.text = f'- Planned vs Actual Hours (Line Chart)'
        chart.chart_title.text_frame.paragraphs[0].font.size = title_font_size  # Adjust chart title font size

        # Adjust font size for labels and category axis
        for series in chart.series:
            for point in series.points:
                point.data_label.font.size = label_font_size

        category_axis = chart.category_axis
        category_axis.tick_labels.font.size = category_axis_font_size

        value_axis = chart.value_axis
        value_axis.tick_labels.font.size = value_axis_font_size

        # Pie Chart
        pie_chart_data = CategoryChartData()
        pie_chart_data.categories = df[df.columns[0]]
        pie_chart_data.add_series('Planned Hours', df['Planned Hours'])

        chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, x3, y3, cx3, cy3, pie_chart_data).chart

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False
        chart.legend.font.size = legend_font_size  # Adjust legend font size
        chart.has_title = True
        chart.chart_title.text_frame.text = f'- Planned Hours Distribution (Pie Chart)'
        chart.chart_title.text_frame.paragraphs[0].font.size = title_font_size  # Adjust chart title font size

        # Adjust font size for labels
        for series in chart.series:
            for point in series.points:
                point.data_label.font.size = label_font_size

        # Area Chart
        area_chart_data = CategoryChartData()
        area_chart_data.categories = df[df.columns[0]]
        area_chart_data.add_series('Planned Hours', df['Planned Hours'])
        area_chart_data.add_series('Actual Hours', df['Actual Hours'])

        chart = slide.shapes.add_chart(XL_CHART_TYPE.AREA, x4, y4, cx4, cy4, area_chart_data).chart

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = legend_font_size  # Adjust legend font size
        chart.has_title = True
        chart.chart_title.text_frame.text = f'- Planned vs Actual Hours (Area Chart)'
        chart.chart_title.text_frame.paragraphs[0].font.size = title_font_size  # Adjust chart title font size

        # Adjust font size for labels and category axis
        for series in chart.series:
            for point in series.points:
                point.data_label.font.size = label_font_size

        category_axis = chart.category_axis
        category_axis.tick_labels.font.size = category_axis_font_size

        value_axis = chart.value_axis
        value_axis.tick_labels.font.size = value_axis_font_size

    except Exception as e:
        print(f"Error adding charts to presentation for {filename}: {e}")


def read_excel_files(input_folder):
    excel_files = [f for f in os.listdir(input_folder) if f.endswith('.xlsx') or f.endswith('.xls')]

    prs = Presentation()

    for filename in excel_files:
        file_path = os.path.join(input_folder, filename)
        try:
            # Read the Excel file into a DataFrame
            df = pd.read_excel(file_path, sheet_name=0, usecols='A:C')

            # Drop rows with any missing values
            df.dropna(inplace=True)

            # Add charts to presentation with custom font sizes
            base_filename = os.path.splitext(filename)[0]
            add_charts_to_presentation(prs, df, base_filename, title_font_size=Pt(15), label_font_size=Pt(10),
                                       category_axis_font_size=Pt(10), value_axis_font_size=Pt(10), legend_font_size=Pt(10))

            print(f"Successfully processed {filename}")

        except Exception as e:
            print(f"Error reading {filename}: {e}")

    # Save the PowerPoint presentation outside the input folder
    prs.save(output_pptx)
    print(f"PowerPoint presentation saved to {output_pptx}")


if __name__ == "__main__":
    read_excel_files(input_folder)'''
# Changes:
# Disabling Gridlines: Added the code to disable gridlines for value_axis in each chart (except the pie chart, which doesn't have gridlines).
'''import os
import pandas as pd
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_TICK_MARK
from pptx.dml.color import RGBColor

# Define the input and output directories
input_folder = 'input'
output_pptx = 'WSR_MultipleWorksheet.pptx'


def add_charts_to_presentation(prs, df, filename, title_font_size=Pt(12), label_font_size=Pt(8),
                               category_axis_font_size=Pt(8), value_axis_font_size=Pt(8), legend_font_size=Pt(8)):
    try:
        required_columns = ['Planned Hours', 'Actual Hours']

        if not all(col in df.columns for col in required_columns):
            print(f"Error: One or more required columns are missing in {filename}. Skipping this file.")
            return

        # Title layout
        slide_layout = prs.slide_layouts[5]  # Title and Content layout

        # Create a new slide for the charts
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = f'{filename} - Charts'
        title.text_frame.paragraphs[0].font.size = title_font_size  # Adjust title font size here

        # Chart size and position
        gap = Cm(0.5)  # Horizontal gap between charts

        # Chart size and position # (HP,VP,  W,H)
        x1, y1, cx1, cy1 = Cm(1.5), Cm(4), Cm(9.9), Cm(6)
        x2, y2, cx2, cy2 = x1 + cx1 + gap, Cm(4), Cm(9.9), Cm(6)
        x3, y3, cx3, cy3 = Cm(1.5), y1 + cy1 + gap, Cm(9.9), Cm(6)
        x4, y4, cx4, cy4 = x3 + cx3 + gap, y3, Cm(9.9), Cm(6)

        # Bar Chart
        bar_chart_data = CategoryChartData()
        bar_chart_data.categories = df[df.columns[0]]
        bar_chart_data.add_series('Planned Hours', df['Planned Hours'])
        bar_chart_data.add_series('Actual Hours', df['Actual Hours'])

        chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x1, y1, cx1, cy1, bar_chart_data).chart

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = legend_font_size  # Adjust legend font size
        chart.has_title = True
        chart.chart_title.text_frame.text = f'- Planned vs Actual Hours (Bar Chart)'
        chart.chart_title.text_frame.paragraphs[0].font.size = title_font_size  # Adjust chart title font size

        # Disable gridlines
        chart.value_axis.major_gridlines.format.line.fill.solid()  # Ensure the line exists
        chart.value_axis.major_gridlines.format.line.fill.background()  # Set fill to background to hide it

        # Adjust font size for labels and category axis
        for series in chart.series:
            for point in series.points:
                point.data_label.font.size = label_font_size

        category_axis = chart.category_axis
        category_axis.tick_labels.font.size = category_axis_font_size

        value_axis = chart.value_axis
        value_axis.tick_labels.font.size = value_axis_font_size

        # Line Chart
        line_chart_data = CategoryChartData()
        line_chart_data.categories = df[df.columns[0]]
        line_chart_data.add_series('Planned Hours', df['Planned Hours'])
        line_chart_data.add_series('Actual Hours', df['Actual Hours'])

        chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, x2, y2, cx2, cy2, line_chart_data).chart

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = legend_font_size  # Adjust legend font size
        chart.has_title = True
        chart.chart_title.text_frame.text = f'- Planned vs Actual Hours (Line Chart)'
        chart.chart_title.text_frame.paragraphs[0].font.size = title_font_size  # Adjust chart title font size

        # Disable gridlines
        chart.value_axis.major_gridlines.format.line.fill.solid()  # Ensure the line exists
        chart.value_axis.major_gridlines.format.line.fill.background()  # Set fill to background to hide it

        # Adjust font size for labels and category axis
        for series in chart.series:
            for point in series.points:
                point.data_label.font.size = label_font_size

        category_axis = chart.category_axis
        category_axis.tick_labels.font.size = category_axis_font_size

        value_axis = chart.value_axis
        value_axis.tick_labels.font.size = value_axis_font_size

        # Pie Chart
        pie_chart_data = CategoryChartData()
        pie_chart_data.categories = df[df.columns[0]]
        pie_chart_data.add_series('Planned Hours', df['Planned Hours'])

        chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, x3, y3, cx3, cy3, pie_chart_data).chart

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False
        chart.legend.font.size = legend_font_size  # Adjust legend font size
        chart.has_title = True
        chart.chart_title.text_frame.text = f'- Planned Hours Distribution (Pie Chart)'
        chart.chart_title.text_frame.paragraphs[0].font.size = title_font_size  # Adjust chart title font size

        # Pie charts do not have gridlines to disable

        # Adjust font size for labels
        for series in chart.series:
            for point in series.points:
                point.data_label.font.size = label_font_size

        # Area Chart
        area_chart_data = CategoryChartData()
        area_chart_data.categories = df[df.columns[0]]
        area_chart_data.add_series('Planned Hours', df['Planned Hours'])
        area_chart_data.add_series('Actual Hours', df['Actual Hours'])

        chart = slide.shapes.add_chart(XL_CHART_TYPE.AREA, x4, y4, cx4, cy4, area_chart_data).chart

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = legend_font_size  # Adjust legend font size
        chart.has_title = True
        chart.chart_title.text_frame.text = f'- Planned vs Actual Hours (Area Chart)'
        chart.chart_title.text_frame.paragraphs[0].font.size = title_font_size  # Adjust chart title font size

        # Disable gridlines
        chart.value_axis.major_gridlines.format.line.fill.solid()  # Ensure the line exists
        chart.value_axis.major_gridlines.format.line.fill.background()  # Set fill to background to hide it

        # Adjust font size for labels and category axis
        for series in chart.series:
            for point in series.points:
                point.data_label.font.size = label_font_size

        category_axis = chart.category_axis
        category_axis.tick_labels.font.size = category_axis_font_size

        value_axis = chart.value_axis
        value_axis.tick_labels.font.size = value_axis_font_size

    except Exception as e:
        print(f"Error adding charts to presentation for {filename}: {e}")


def read_excel_files(input_folder):
    excel_files = [f for f in os.listdir(input_folder) if f.endswith('.xlsx') or f.endswith('.xls')]

    prs = Presentation()

    for filename in excel_files:
        file_path = os.path.join(input_folder, filename)
        try:
            # Read the Excel file into a DataFrame
            df = pd.read_excel(file_path, sheet_name=0, usecols='A:C')

            # Drop rows with any missing values
            df.dropna(inplace=True)

            # Add charts to presentation with custom font sizes
            base_filename = os.path.splitext(filename)[0]
            add_charts_to_presentation(prs, df, base_filename, title_font_size=Pt(15), label_font_size=Pt(10),
                                       category_axis_font_size=Pt(10), value_axis_font_size=Pt(10),
                                       legend_font_size=Pt(10))

            print(f"Successfully processed {filename}")

        except Exception as e:
            print(f"Error reading {filename}: {e}")

    # Save the PowerPoint presentation outside the input folder
    prs.save(output_pptx)
    print(f"PowerPoint presentation saved to {output_pptx}")


if __name__ == "__main__":
    read_excel_files(input_folder)'''
# To adjust the series overlap for the primary axis in the charts, you can set the overlap property of the chart.series object. Below is the updated code with the series overlap adjusted for the bar chart, line chart, and area chart:
import os
import pandas as pd
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_TICK_MARK
from pptx.dml.color import RGBColor

# Define the input and output directories
input_folder = 'input'
output_pptx = 'WSR_MultipleWorksheet.pptx'

def add_charts_to_presentation(prs, df, filename, title_font_size=Pt(12), label_font_size=Pt(8),
                               category_axis_font_size=Pt(8), value_axis_font_size=Pt(8), legend_font_size=Pt(8), series_overlap=-27):
    try:
        required_columns = ['Planned Hours', 'Actual Hours']

        if not all(col in df.columns for col in required_columns):
            print(f"Error: One or more required columns are missing in {filename}. Skipping this file.")
            return

        # Title layout
        slide_layout = prs.slide_layouts[5]  # Title and Content layout

        # Create a new slide for the charts
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = f'{filename} - Charts'
        title.text_frame.paragraphs[0].font.size = title_font_size  # Adjust title font size here

        # Chart size and position
        gap = Cm(0.5)  # Horizontal gap between charts

        # Chart size and position
        x1, y1, cx1, cy1 = Cm(1.5), Cm(4), Cm(9.9), Cm(6)
        x2, y2, cx2, cy2 = x1 + cx1 + gap, Cm(4), Cm(9.9), Cm(6)
        x3, y3, cx3, cy3 = Cm(1.5), y1 + cy1 + gap, Cm(9.9), Cm(6)
        x4, y4, cx4, cy4 = x3 + cx3 + gap, y3, Cm(9.9), Cm(6)

        # Bar Chart
        bar_chart_data = CategoryChartData()
        bar_chart_data.categories = df[df.columns[0]]
        bar_chart_data.add_series('Planned Hours', df['Planned Hours'])
        bar_chart_data.add_series('Actual Hours', df['Actual Hours'])

        chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x1, y1, cx1, cy1, bar_chart_data).chart

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = legend_font_size  # Adjust legend font size
        chart.has_title = True
        chart.chart_title.text_frame.text = f'- Planned vs Actual Hours (Bar Chart)'
        chart.chart_title.text_frame.paragraphs[0].font.size = title_font_size  # Adjust chart title font size

        # Disable gridlines
        chart.value_axis.major_gridlines.format.line.fill.solid()  # Ensure the line exists
        chart.value_axis.major_gridlines.format.line.fill.background()  # Set fill to background to hide it

        # Adjust font size for labels and category axis
        for series in chart.series:
            for point in series.points:
                point.data_label.font.size = label_font_size

        category_axis = chart.category_axis
        category_axis.tick_labels.font.size = category_axis_font_size

        value_axis = chart.value_axis
        value_axis.tick_labels.font.size = value_axis_font_size

        # Adjust series overlap for bar chart
        chart.plots[0].overlap = series_overlap

        # Line Chart
        line_chart_data = CategoryChartData()
        line_chart_data.categories = df[df.columns[0]]
        line_chart_data.add_series('Planned Hours', df['Planned Hours'])
        line_chart_data.add_series('Actual Hours', df['Actual Hours'])

        chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, x2, y2, cx2, cy2, line_chart_data).chart

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = legend_font_size  # Adjust legend font size
        chart.has_title = True
        chart.chart_title.text_frame.text = f'- Planned vs Actual Hours (Line Chart)'
        chart.chart_title.text_frame.paragraphs[0].font.size = title_font_size  # Adjust chart title font size

        # Disable gridlines
        chart.value_axis.major_gridlines.format.line.fill.solid()  # Ensure the line exists
        chart.value_axis.major_gridlines.format.line.fill.background()  # Set fill to background to hide it

        # Adjust font size for labels and category axis
        for series in chart.series:
            for point in series.points:
                point.data_label.font.size = label_font_size

        category_axis = chart.category_axis
        category_axis.tick_labels.font.size = category_axis_font_size

        value_axis = chart.value_axis
        value_axis.tick_labels.font.size = value_axis_font_size

        # Pie Chart
        pie_chart_data = CategoryChartData()
        pie_chart_data.categories = df[df.columns[0]]
        pie_chart_data.add_series('Planned Hours', df['Planned Hours'])

        chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, x3, y3, cx3, cy3, pie_chart_data).chart

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False
        chart.legend.font.size = legend_font_size  # Adjust legend font size
        chart.has_title = True
        chart.chart_title.text_frame.text = f'- Planned Hours Distribution (Pie Chart)'
        chart.chart_title.text_frame.paragraphs[0].font.size = title_font_size  # Adjust chart title font size

        # Pie charts do not have gridlines to disable

        # Adjust font size for labels
        for series in chart.series:
            for point in series.points:
                point.data_label.font.size = label_font_size

        # Area Chart
        area_chart_data = CategoryChartData()
        area_chart_data.categories = df[df.columns[0]]
        area_chart_data.add_series('Planned Hours', df['Planned Hours'])
        area_chart_data.add_series('Actual Hours', df['Actual Hours'])

        chart = slide.shapes.add_chart(XL_CHART_TYPE.AREA, x4, y4, cx4, cy4, area_chart_data).chart

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = legend_font_size  # Adjust legend font size
        chart.has_title = True
        chart.chart_title.text_frame.text = f'- Planned vs Actual Hours (Area Chart)'
        chart.chart_title.text_frame.paragraphs[0].font.size = title_font_size  # Adjust chart title font size

        # Disable gridlines
        chart.value_axis.major_gridlines.format.line.fill.solid()  # Ensure the line exists
        chart.value_axis.major_gridlines.format.line.fill.background()  # Set fill to background to hide it

        # Adjust font size for labels and category axis
        for series in chart.series:
            for point in series.points:
                point.data_label.font.size = label_font_size

        category_axis = chart.category_axis
        category_axis.tick_labels.font.size = category_axis_font_size

        value_axis = chart.value_axis
        value_axis.tick_labels.font.size = value_axis_font_size

        # Adjust series overlap for area chart
        chart.plots[0].overlap = series_overlap

    except Exception as e:
        print(f"Error adding charts to presentation for {filename}: {e}")


def read_excel_files(input_folder):
    excel_files = [f for f in os.listdir(input_folder) if f.endswith('.xlsx') or f.endswith('.xls')]

    prs = Presentation()

    for filename in excel_files:
        file_path = os.path.join(input_folder, filename)
        try:
            # Read the Excel file into a DataFrame
            df = pd.read_excel(file_path, sheet_name=0, usecols='A:C')

            # Drop rows with any missing values
            df.dropna(inplace=True)

            # Add charts to presentation with custom font sizes
            base_filename = os.path.splitext(filename)[0]
            add_charts_to_presentation(prs, df, base_filename, title_font_size=Pt(15), label_font_size=Pt(10),
                                       category_axis_font_size=Pt(10), value_axis_font_size=Pt(10), legend_font_size=Pt(10), series_overlap=-27)

            print(f"Successfully processed {filename}")

        except Exception as e:
            print(f"Error reading {filename}: {e}")

    # Save the PowerPoint presentation outside the input folder
    prs.save(output_pptx)
    print(f"PowerPoint presentation saved to {output_pptx}")


if __name__ == "__main__":
    read_excel_files(input_folder)



