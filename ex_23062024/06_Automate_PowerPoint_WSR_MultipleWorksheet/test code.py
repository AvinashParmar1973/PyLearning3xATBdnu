# o adjust the series overlap for the secondary axis in the line chart and area chart, you need to first add a secondary axis and then adjust the overlap for the series plotted on that axis. Here is the updated code with these changes:
import os
import pandas as pd
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_TICK_MARK
from pptx.dml.color import RGBColor

# Define the input and output directories
input_folder = 'input'
output_pptx = 'WSR_MultipleWorksheet.pptx'

def add_charts_to_presentation(prs, df, filename, title_font_size=Pt(12), label_font_size=Pt(8),
                               category_axis_font_size=Pt(8), value_axis_font_size=Pt(8), legend_font_size=Pt(8), series_overlap=-27):
    try:
        required_columns = ['Planned Hours', 'Actual Hours']

        if not all(col in df.columns for col in required_columns):
            print(f"Error: One or more required columns are missing in {filename}. Skipping this file.")
            return

        # Title layout
        slide_layout = prs.slide_layouts[5]  # Title and Content layout

        # Create a new slide for the charts
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = f'{filename} - Charts'
        title.text_frame.paragraphs[0].font.size = title_font_size  # Adjust title font size here

        # Chart size and position
        gap = Cm(0.5)  # Horizontal gap between charts

        x1, y1, cx1, cy1 = Cm(0.5), Cm(4), Cm(9.5), Cm(6)
        x2, y2, cx2, cy2 = x1 + cx1 + gap, Cm(4), Cm(9.5), Cm(6)
        x3, y3, cx3, cy3 = Cm(0.5), y1 + cy1 + gap, Cm(9.5), Cm(6)
        x4, y4, cx4, cy4 = x3 + cx3 + gap, y3, Cm(9.5), Cm(6)

        # Bar Chart
        bar_chart_data = CategoryChartData()
        bar_chart_data.categories = df[df.columns[0]]
        bar_chart_data.add_series('Planned Hours', df['Planned Hours'])
        bar_chart_data.add_series('Actual Hours', df['Actual Hours'])

        chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x1, y1, cx1, cy1, bar_chart_data).chart

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = legend_font_size  # Adjust legend font size
        chart.has_title = True
        chart.chart_title.text_frame.text = f'- Planned vs Actual Hours (Bar Chart)'
        chart.chart_title.text_frame.paragraphs[0].font.size = title_font_size  # Adjust chart title font size

        # Disable gridlines
        chart.value_axis.major_gridlines.format.line.fill.solid()  # Ensure the line exists
        chart.value_axis.major_gridlines.format.line.fill.background()  # Set fill to background to hide it

        # Adjust font size for labels and category axis
        for series in chart.series:
            for point in series.points:
                point.data_label.font.size = label_font_size

        category_axis = chart.category_axis
        category_axis.tick_labels.font.size = category_axis_font_size

        value_axis = chart.value_axis
        value_axis.tick_labels.font.size = value_axis_font_size

        # Adjust series overlap for bar chart
        chart.plots[0].overlap = series_overlap

        # Line Chart
        line_chart_data = CategoryChartData()
        line_chart_data.categories = df[df.columns[0]]
        line_chart_data.add_series('Planned Hours', df['Planned Hours'])
        line_chart_data.add_series('Actual Hours', df['Actual Hours'])

        chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, x2, y2, cx2, cy2, line_chart_data).chart

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = legend_font_size  # Adjust legend font size
        chart.has_title = True
        chart.chart_title.text_frame.text = f'- Planned vs Actual Hours (Line Chart)'
        chart.chart_title.text_frame.paragraphs[0].font.size = title_font_size  # Adjust chart title font size

        # Disable gridlines
        chart.value_axis.major_gridlines.format.line.fill.solid()  # Ensure the line exists
        chart.value_axis.major_gridlines.format.line.fill.background()  # Set fill to background to hide it

        # Adjust font size for labels and category axis
        for series in chart.series:
            for point in series.points:
                point.data_label.font.size = label_font_size

        category_axis = chart.category_axis
        category_axis.tick_labels.font.size = category_axis_font_size

        value_axis = chart.value_axis
        value_axis.tick_labels.font.size = value_axis_font_size

        # Add secondary value axis and adjust series overlap for line chart
        secondary_value_axis = chart.value_axis
        chart.has_secondary_value_axis = True
        secondary_value_axis = chart.value_axis
        chart.plots[0].overlap = series_overlap

        # Pie Chart
        pie_chart_data = CategoryChartData()
        pie_chart_data.categories = df[df.columns[0]]
        pie_chart_data.add_series('Planned Hours', df['Planned Hours'])

        chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, x3, y3, cx3, cy3, pie_chart_data).chart

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False
        chart.legend.font.size = legend_font_size  # Adjust legend font size
        chart.has_title = True
        chart.chart_title.text_frame.text = f'- Planned Hours Distribution (Pie Chart)'
        chart.chart_title.text_frame.paragraphs[0].font.size = title_font_size  # Adjust chart title font size

        # Pie charts do not have gridlines to disable

        # Adjust font size for labels
        for series in chart.series:
            for point in series.points:
                point.data_label.font.size = label_font_size

        # Area Chart
        area_chart_data = CategoryChartData()
        area_chart_data.categories = df[df.columns[0]]
        area_chart_data.add_series('Planned Hours', df['Planned Hours'])
        area_chart_data.add_series('Actual Hours', df['Actual Hours'])

        chart = slide.shapes.add_chart(XL_CHART_TYPE.AREA, x4, y4, cx4, cy4, area_chart_data).chart

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = legend_font_size  # Adjust legend font size
        chart.has_title = True
        chart.chart_title.text_frame.text = f'- Planned vs Actual Hours (Area Chart)'
        chart.chart_title.text_frame.paragraphs[0].font.size = title_font_size  # Adjust chart title font size

        # Disable gridlines
        chart.value_axis.major_gridlines.format.line.fill.solid()  # Ensure the line exists
        chart.value_axis.major_gridlines.format.line.fill.background()  # Set fill to background to hide it

        # Adjust font size for labels and category axis
        for series in chart.series:
            for point in series.points:
                point.data_label.font.size = label_font_size

        category_axis = chart.category_axis
        category_axis.tick_labels.font.size = category_axis_font_size

        value_axis = chart.value_axis
        value_axis.tick_labels.font.size = value_axis_font_size

        # Add secondary value axis and adjust series overlap for area chart
        secondary_value_axis = chart.value_axis
        chart.has_secondary_value_axis = True
        secondary_value_axis = chart.value_axis
        chart.plots[0].overlap = series_overlap

    except Exception as e:
        print(f"Error adding charts to presentation for {filename}: {e}")


def read_excel_files(input_folder):
    excel_files = [f for f in os.listdir(input_folder) if f.endswith('.xlsx') or f.endswith('.xls')]

    prs = Presentation()

    for filename in excel_files:
        file_path = os.path.join(input_folder, filename)
        try:
            # Read the Excel file into a DataFrame
            df = pd.read_excel(file_path, sheet_name=0, usecols='A:C')

            # Drop rows with any missing values
            df.dropna(inplace=True)

            # Add charts to presentation with custom font sizes
            base_filename = os.path.splitext(filename)[0]
            add_charts_to_presentation(prs, df, base_filename, title_font_size=Pt(15), label_font_size=Pt(10),
                                       category_axis_font_size=Pt(10), value_axis_font_size=Pt(10), legend_font_size=Pt(10), series_overlap=-27)

            print(f"Successfully processed {filename}")

        except Exception as e:
            print(f"Error reading {filename}: {e}")

    # Save the PowerPoint presentation outside the input folder
    prs.save(output_pptx)
    print(f"PowerPoint presentation saved to {output_pptx}")


if __name__ == "__main__":
    read_excel_files(input_folder)
